import os
import sys

# Tự động kiểm tra và cài đặt thư viện python-pptx nếu chưa có
try:
    import pptx
except ImportError:
    print("Không tìm thấy thư viện 'python-pptx'. Đang tiến hành cài đặt tự động...")
    import subprocess
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
        print("Đã cài đặt 'python-pptx' thành công!\n")
    except Exception as e:
        print(f"Cài đặt thất bại. Vui lòng cài đặt thủ công bằng lệnh: pip install python-pptx")
        print(f"Chi tiết lỗi: {e}")
        sys.exit(1)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Khởi tạo bản trình chiếu
prs = Presentation()

# Đặt kích thước slide chuẩn 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Bảng phối màu (Dark Theme chuyên nghiệp)
BG_COLOR = RGBColor(15, 23, 42)        # Xanh đen Navy tối giản
TEXT_TITLE = RGBColor(250, 204, 21)    # Vàng rực (Warm Gold) cho Tiêu đề
TEXT_BODY = RGBColor(226, 232, 240)     # Trắng bạc (Light Silver) cho nội dung
TEXT_HIGHLIGHT = RGBColor(6, 182, 212) # Cyan (Xanh neon) cho điểm nhấn
CARD_BG = RGBColor(30, 41, 59)          # Màu nền cho hộp nội dung (Slate)

def set_slide_background(slide):
    """Tạo hình nền tối màu cho toàn bộ slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, title_text):
    """Thêm tiêu đề chuẩn cho slide nội dung"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    return title_box

# ==========================================
# SLIDE 1: SLIDE TIÊU ĐỀ
# ==========================================
slide_layout = prs.slide_layouts[6] # Trống hoàn toàn
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1)

# Thêm hộp văn bản Tiêu đề
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

# Tiêu đề chính
p1 = tf.paragraphs[0]
p1.text = "BÀI 1: CẤU TRÚC CỦA CHẤT. SỰ CHUYỂN THỂ"
p1.alignment = PP_ALIGN.CENTER
p1.font.name = 'Arial'
p1.font.size = Pt(40)
p1.font.bold = True
p1.font.color.rgb = TEXT_TITLE
p1.space_after = Pt(20)

# Tiêu đề phụ
p2 = tf.add_paragraph()
p2.text = "Chương I: Vật lí nhiệt  |  Vật lí 12 (GDPT 2018)"
p2.alignment = PP_ALIGN.CENTER
p2.font.name = 'Arial'
p2.font.size = Pt(20)
p2.font.color.rgb = TEXT_BODY
p2.space_after = Pt(40)

# Giáo viên
p3 = tf.add_paragraph()
p3.text = "Giáo viên thực hiện: [Tên của bạn]"
p3.alignment = PP_ALIGN.CENTER
p3.font.name = 'Arial'
p3.font.size = Pt(16)
p3.font.italic = True
p3.font.color.rgb = TEXT_HIGHLIGHT

# ==========================================
# SLIDE 2: KHỞI ĐỘNG - HIỆN TƯỢNG THỰC TẾ
# ==========================================
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2)
add_header(slide2, "HIỆN TƯỢNG QUAN SÁT THỰC TẾ")

# Tạo 3 cột nội dung (3 Cards) đại diện cho 3 trạng thái chuyển đổi
cols_data = [
    ("1. BĂNG TAN", "Khối đá rắn góc cạnh tan chảy dưới nắng ấm thành nước lỏng mềm mại."),
    ("2. SƯƠNG BỐC HƠI", "Giọt sương sớm đọng trên lá cây bốc hơi thành dải sương mù vô hình."),
    ("3. BĂNG KHÔ KHÓI", "Băng khô thăng hoa trực tiếp thành khí cuộn cuộn mà không chảy lỏng.")
]

left_margin = Inches(0.8)
col_width = Inches(3.6)
gap = Inches(0.5)

for i, (title, desc) in enumerate(cols_data):
    x = left_margin + i * (col_width + gap)
    # Vẽ hộp chứa (Card)
    shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), col_width, Inches(3.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = TEXT_HIGHLIGHT
    shape.line.width = Pt(1.5)
    
    # Ghi nội dung vào hộp
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = 'Arial'
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_BODY
    p2.line_spacing = 1.2

# Câu hỏi gợi mở ở dưới cùng
question_box = slide2.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.733), Inches(1.0))
tf = question_box.text_frame
tf.word_wrap = True
p_q = tf.paragraphs[0]
p_q.text = "CHÚNG TA SUY NGẪM: Bản chất đằng sau các quá trình này là gì? Chúng được cấu tạo thế nào?"
p_q.alignment = PP_ALIGN.CENTER
p_q.font.name = 'Arial'
p_q.font.size = Pt(18)
p_q.font.bold = True
p_q.font.color.rgb = TEXT_TITLE

# ==========================================
# SLIDE 3: BẰNG CHỨNG THỰC NGHIỆM - CHUYỂN ĐỘNG BROWN
# ==========================================
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3)
add_header(slide3, "BẰNG CHỨNG THỰC NGHIỆM VỀ NGUYÊN TỬ")

# Trái: Khung chứa hình minh họa
pic_box = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.0), Inches(4.5))
pic_box.fill.solid()
pic_box.fill.fore_color.rgb = CARD_BG
pic_box.line.color.rgb = TEXT_BODY
tf_pic = pic_box.text_frame
p_pic = tf_pic.paragraphs[0]
p_pic.text = "[MÔ PHỎNG CHUYỂN ĐỘNG BROWN]\nHạt phấn hoa bị va đập bởi các phân tử nước vô hình chuyển động liên tục không ngừng"
p_pic.alignment = PP_ALIGN.CENTER
p_pic.font.name = 'Arial'
p_pic.font.size = Pt(16)
p_pic.font.color.rgb = TEXT_HIGHLIGHT

# Phải: Thông tin văn bản
text_box = slide3.shapes.add_textbox(Inches(6.2), Inches(1.5), Inches(6.333), Inches(4.5))
tf = text_box.text_frame
tf.word_wrap = True

bullets = [
    ("Thí nghiệm Robert Brown (1827):", "Quan sát thấy các hạt phấn hoa lơ lửng trong nước chuyển động hỗn loạn không ngừng về mọi phía."),
    ("Giải thích của Albert Einstein (1905):", "Chuyển động hỗn loạn là do các phân tử nước siêu nhỏ, vô hình trước mắt thường, chuyển động nhiệt va chạm liên tục từ mọi hướng."),
    ("Ý nghĩa khoa học cốt lõi:", "Khẳng định gián tiếp sự tồn tại thực sự của nguyên tử, phân tử và chứng minh chúng chuyển động không ngừng.")
]

for i, (title, body) in enumerate(bullets):
    p_t = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p_t.text = title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_TITLE
    if i > 0: p_t.space_before = Pt(14)
    p_t.space_after = Pt(4)
    
    p_b = tf.add_paragraph()
    p_b.text = body
    p_b.font.name = 'Arial'
    p_b.font.size = Pt(16)
    p_b.font.color.rgb = TEXT_BODY

# ==========================================
# SLIDE 4: THUYẾT ĐỘNG HỌC PHÂN TỬ CHẤT
# ==========================================
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4)
add_header(slide4, "THUYẾT ĐỘNG HỌC PHÂN TỬ CHẤT")

# 3 hộp nội dung đặt dọc
rows_data = [
    ("1. CẤU TẠO HẠT VÀ KHOẢNG CÁCH", "Các chất được cấu tạo từ các hạt riêng biệt gọi là phân tử (hoặc nguyên tử). Giữa các phân tử luôn có khoảng cách nhất định."),
    ("2. CHUYỂN ĐỘNG NHIỆT KHÔNG NGỪNG", "Các phân tử chuyển động hỗn loạn không ngừng. Nhiệt độ của vật càng cao, tốc độ chuyển động nhiệt trung bình của các phân tử càng lớn."),
    ("3. LỰC TƯƠNG TÁC PHÂN TỬ SONG HÀNH", "Giữa các phân tử đồng thời tồn tại cả lực hút và lực đẩy phân tử, giữ liên kết hoặc đẩy các hạt xa nhau tùy khoảng cách.")
]

for i, (title, desc) in enumerate(rows_data):
    y = Inches(1.5) + i * Inches(1.8)
    shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.733), Inches(1.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = TEXT_HIGHLIGHT
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.space_after = Pt(4)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = 'Arial'
    p2.font.size = Pt(15)
    p2.font.color.rgb = TEXT_BODY

# ==========================================
# SLIDE 5: BẢN CHẤT LỰC TƯƠNG TÁC PHÂN TỬ
# ==========================================
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5)
add_header(slide5, "BẢN CHẤT LỰC TƯƠNG TÁC PHÂN TỬ")

# Trái: Đồ họa lò xo
box_left = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.4), Inches(4.8))
box_left.fill.solid()
box_left.fill.fore_color.rgb = CARD_BG
box_left.line.color.rgb = TEXT_HIGHLIGHT
tf_bl = box_left.text_frame
tf_bl.word_wrap = True
p_bl = tf_bl.paragraphs[0]
p_bl.text = "[MÔ HÌNH LÒ XO ĐÀN HỒI 3D]\n\nXem phân tử như các hạt tròn nối với nhau bằng lò xo.\n- Co lò xo: Lực đẩy cực mạnh xuất hiện.\n- Giãn lò xo: Lực hút kéo hai hạt lại gần."
p_bl.alignment = PP_ALIGN.CENTER
p_bl.font.name = 'Arial'
p_bl.font.size = Pt(16)
p_bl.font.color.rgb = TEXT_BODY

# Phải: Đồ thị lực tương tác F(r)
box_right = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.8))
box_right.fill.solid()
box_right.fill.fore_color.rgb = CARD_BG
box_right.line.color.rgb = TEXT_BODY
tf_br = box_right.text_frame
tf_br.word_wrap = True
p_br = tf_br.paragraphs[0]
p_br.text = "[ĐỒ THỊ LỰC TƯƠNG TÁC F(r)]\n\n- Trục dọc hiển thị độ lớn lực tương tác.\n- Điểm cắt trục hoành tại r0 là nơi F = 0.\n- Vùng r < r0 biểu diễn lực đẩy (dương).\n- Vùng r > r0 biểu diễn lực hút (âm)."
p_br.alignment = PP_ALIGN.CENTER
p_br.font.name = 'Arial'
p_br.font.size = Pt(16)
p_br.font.color.rgb = TEXT_HIGHLIGHT

# ==========================================
# SLIDE 6: LỰC TƯƠNG TÁC THEO KHOẢNG CÁCH
# ==========================================
slide6 = prs.slides.add_slide(slide_layout)
set_slide_background(slide6)
add_header(slide6, "LỰC TƯƠNG TÁC VÀ KHOẢNG CÁCH r")

# Tạo 4 khối giải thích 4 trường hợp khoảng cách r
cases = [
    ("1. r = r0 (Cân bằng)", "Khoảng cách ổn định. Lực hút bằng lực đẩy, lực tổng hợp F = 0. Hai phân tử dao động tại chỗ."),
    ("2. r < r0 (Ép lại gần)", "Hai phân tử quá gần, lò xo co. Lực đẩy lớn hơn lực hút, lực tổng hợp là lực đẩy chống lại nén."),
    ("3. r > r0 (Kéo ra xa)", "Hai phân tử xa nhau, lò xo giãn. Lực hút lớn hơn lực đẩy, lực tổng hợp là lực hút kéo xích lại."),
    ("4. r >> r0 (Quá xa nhau)", "Lực tương tác giảm nhanh chóng về 0. Các phân tử không liên kết (coi như chất khí lý tưởng).")
]

for i, (title, desc) in enumerate(cases):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.0)
    y = Inches(1.5) + row * Inches(2.6)
    
    shape = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = TEXT_HIGHLIGHT
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.space_after = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = 'Arial'
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_BODY
    p2.line_spacing = 1.15

# ==========================================
# SLIDE 7: MÔ HÌNH CẤU TRÚC THỂ RẮN
# ==========================================
slide7 = prs.slides.add_slide(slide_layout)
set_slide_background(slide7)
add_header(slide7, "MÔ HÌNH CẤU TRÚC THỂ RẮN (SOLID)")

box_l = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.4), Inches(4.8))
box_l.fill.solid()
box_l.fill.fore_color.rgb = CARD_BG
box_l.line.color.rgb = TEXT_BODY
tf_bl = box_l.text_frame
p_bl = tf_bl.paragraphs[0]
p_bl.text = "[MẠNG TINH THỂ RẮN 3D]\n\nCác hạt khăng khít sắp xếp tuần hoàn có trật tự trong không gian 3 chiều và dao động khẽ."
p_bl.alignment = PP_ALIGN.CENTER
p_bl.font.name = 'Arial'
p_bl.font.size = Pt(16)
p_bl.font.color.rgb = TEXT_HIGHLIGHT

box_r = slide7.shapes.add_textbox(Inches(6.6), Inches(1.5), Inches(5.9), Inches(4.8))
tf = box_r.text_frame
tf.word_wrap = True

details = [
    ("Khoảng cách phân tử:", "Rất nhỏ, xấp xỉ bằng kích thước phân tử (r ≈ r0)."),
    ("Lực tương tác phân tử:", "Rất mạnh, khóa chặt các phân tử tại các điểm nút mạng."),
    ("Chuyển động phân tử:", "Chỉ dao động nhiệt nhỏ tại chỗ quanh vị trí cân bằng cố định."),
    ("Đặc điểm vĩ mô:", "Có hình dạng và thể tích riêng xác định, hầu như rất khó nén.")
]

for i, (title, desc) in enumerate(details):
    p_t = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p_t.text = title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_TITLE
    if i > 0: p_t.space_before = Pt(12)
    p_t.space_after = Pt(2)
    
    p_d = tf.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(15)
    p_d.font.color.rgb = TEXT_BODY

# ==========================================
# SLIDE 8: MÔ HÌNH CẤU TRÚC THỂ LỎNG
# ==========================================
slide8 = prs.slides.add_slide(slide_layout)
set_slide_background(slide8)
add_header(slide8, "MÔ HÌNH CẤU TRÚC THỂ LỎNG (LIQUID)")

box_l = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.4), Inches(4.8))
box_l.fill.solid()
box_l.fill.fore_color.rgb = CARD_BG
box_l.line.color.rgb = TEXT_BODY
tf_bl = box_l.text_frame
p_bl = tf_bl.paragraphs[0]
p_bl.text = "[CẤU TRÚC THỂ LỎNG 3D]\n\nCác hạt phân tử ở gần nhau nhưng xếp lộn xộn, dễ dàng trượt qua nhau tự do dưới lòng chất chứa."
p_bl.alignment = PP_ALIGN.CENTER
p_bl.font.name = 'Arial'
p_bl.font.size = Pt(16)
p_bl.font.color.rgb = TEXT_HIGHLIGHT

box_r = slide8.shapes.add_textbox(Inches(6.6), Inches(1.5), Inches(5.9), Inches(4.8))
tf = box_r.text_frame
tf.word_wrap = True

details = [
    ("Khoảng cách phân tử:", "Nhỏ, lớn hơn một ít so với thể rắn nhưng các phân tử vẫn sát nhau."),
    ("Lực tương tác phân tử:", "Yếu hơn thể rắn, chưa đủ giữ phân tử ở một vị trí cố định."),
    ("Chuyển động phân tử:", "Dao động quanh vị trí cân bằng tạm thời, sau một lúc vị trí này dịch chuyển giúp các hạt trượt tự do."),
    ("Đặc điểm vĩ mô:", "Có thể tích riêng xác định, hình dạng phụ thuộc bình chứa, khó nén.")
]

for i, (title, desc) in enumerate(details):
    p_t = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p_t.text = title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_TITLE
    if i > 0: p_t.space_before = Pt(12)
    p_t.space_after = Pt(2)
    
    p_d = tf.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(15)
    p_d.font.color.rgb = TEXT_BODY

# ==========================================
# SLIDE 9: MÔ HÌNH CẤU TRÚC THỂ KHÍ
# ==========================================
slide9 = prs.slides.add_slide(slide_layout)
set_slide_background(slide9)
add_header(slide9, "MÔ HÌNH CẤU TRÚC THỂ KHÍ (GAS)")

box_l = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.4), Inches(4.8))
box_l.fill.solid()
box_l.fill.fore_color.rgb = CARD_BG
box_l.line.color.rgb = TEXT_BODY
tf_bl = box_l.text_frame
p_bl = tf_bl.paragraphs[0]
p_bl.text = "[CẤU TRÚC CHẤT KHÍ 3D]\n\nCác hạt ở rất xa nhau, chuyển động với tốc độ cao đâm sầm hỗn loạn và tự do trong mọi không gian bình."
p_bl.alignment = PP_ALIGN.CENTER
p_bl.font.name = 'Arial'
p_bl.font.size = Pt(16)
p_bl.font.color.rgb = TEXT_HIGHLIGHT

box_r = slide9.shapes.add_textbox(Inches(6.6), Inches(1.5), Inches(5.9), Inches(4.8))
tf = box_r.text_frame
tf.word_wrap = True

details = [
    ("Khoảng cách phân tử:", "Rất lớn (gấp hàng chục lần đường kính phân tử khí)."),
    ("Lực tương tác phân tử:", "Rất yếu, hầu như có thể bỏ qua ngoại trừ khi chúng va đập."),
    ("Chuyển động phân tử:", "Chuyển động hỗn loạn hoàn toàn tự do. Khi va chạm thì nảy ra hướng khác."),
    ("Đặc điểm vĩ mô:", "Không có hình dạng và thể tích riêng xác định, chiếm toàn bộ bình, cực kỳ dễ nén.")
]

for i, (title, desc) in enumerate(details):
    p_t = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p_t.text = title
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_TITLE
    if i > 0: p_t.space_before = Pt(12)
    p_t.space_after = Pt(2)
    
    p_d = tf.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(15)
    p_d.font.color.rgb = TEXT_BODY

# ==========================================
# SLIDE 10: BẢNG SO SÁNH BA TRẠNG THÁI CỦA CHẤT
# ==========================================
slide10 = prs.slides.add_slide(slide_layout)
set_slide_background(slide10)
add_header(slide10, "SO SÁNH CẤU TRÚC BA THỂ CỦA CHẤT")

# Thiết kế bảng 7 hàng, 4 cột
rows = 7
cols = 4
left = Inches(0.8)
top = Inches(1.4)
width = Inches(11.733)
height = Inches(5.2)

table_shape = slide10.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set độ rộng cột
table.columns[0].width = Inches(2.233)
table.columns[1].width = Inches(3.1)
table.columns[2].width = Inches(3.1)
table.columns[3].width = Inches(3.3)

headers = ["Đặc trưng", "Thể Rắn", "Thể Lỏng", "Thể Khí"]
for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG
    p = cell.text_frame.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.alignment = PP_ALIGN.CENTER

table_data = [
    ["Khoảng cách", "Rất nhỏ (gần sát nhau)", "Nhỏ (gần nhau)", "Rất lớn (ở rất xa nhau)"],
    ["Lực liên kết", "Rất mạnh", "Yếu hơn thể rắn", "Rất yếu (coi như bằng 0)"],
    ["Chuyển động", "Dao động quanh VTCB cố định", "Dao động quanh VTCB tạm thời", "Chuyển động hỗn loạn tự do"],
    ["Hình dạng", "Xác định ổn định", "Phụ thuộc vào bình chứa", "Phụ thuộc vào bình chứa"],
    ["Thể tích", "Xác định ổn định", "Xác định ổn định", "Chiếm toàn bộ bình chứa"],
    ["Khả năng nén", "Rất khó nén", "Khó nén", "Cực kỳ dễ nén"]
]

for row_idx, row_values in enumerate(table_data):
    for col_idx, val in enumerate(row_values):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG if row_idx % 2 == 0 else BG_COLOR
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_BODY
        if col_idx == 0:
            p.font.bold = True
            p.font.color.rgb = TEXT_TITLE

# ==========================================
# SLIDE 11: CÁC QUÁ TRÌNH CHUYỂN THỂ CỦA CHẤT
# ==========================================
slide11 = prs.slides.add_slide(slide_layout)
set_slide_background(slide11)
add_header(slide11, "CÁC QUÁ TRÌNH CHUYỂN THỂ CỦA CHẤT")

# Thiết kế sơ đồ 3 card Rắn, Lỏng, Khí nằm ngang kèm mũi tên trung gian
phases = [
    ("RẮN", ["➔ Nóng chảy ➔ Lỏng", "➔ Thăng hoa ➔ Khí"], CARD_BG),
    ("LỎNG", ["➔ Đông đặc ➔ Rắn", "➔ Hóa hơi ➔ Khí"], CARD_BG),
    ("KHÍ", ["➔ Ngưng tụ ➔ Lỏng", "➔ Ngưng kết ➔ Rắn"], CARD_BG)
]

for i, (title, transitions, bg) in enumerate(phases):
    x = Inches(0.8) + i * Inches(4.1)
    shape = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.6), Inches(3.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = TEXT_HIGHLIGHT
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.4)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)
    
    for trans in transitions:
        p2 = tf.add_paragraph()
        p2.text = trans
        p2.font.name = 'Arial'
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT_BODY
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(10)

# Quy luật năng lượng ở dưới
rules_box = slide11.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.733), Inches(1.0))
tf = rules_box.text_frame
p_rule = tf.paragraphs[0]
p_rule.text = "Thu nhiệt lượng (Năng lượng vào): Nóng chảy, Hóa hơi, Thăng hoa\nTỏa nhiệt lượng (Năng lượng ra): Đông đặc, Ngưng tụ, Ngưng kết"
p_rule.alignment = PP_ALIGN.CENTER
p_rule.font.name = 'Arial'
p_rule.font.size = Pt(16)
p_rule.font.bold = True
p_rule.font.color.rgb = TEXT_TITLE

# ==========================================
# SLIDE 12: GIẢI THÍCH SỰ CHUYỂN THỂ DƯỚI GÓC PHÂN TỬ - NÓNG CHẢY & ĐÔNG ĐẶC
# ==========================================
slide12 = prs.slides.add_slide(slide_layout)
set_slide_background(slide12)
add_header(slide12, "CƠ CHẾ NÓNG CHẢY VÀ ĐÔNG ĐẶC")

# Cột Trái: Sự nóng chảy
box_l = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.8))
box_l.fill.solid()
box_l.fill.fore_color.rgb = CARD_BG
box_l.line.color.rgb = TEXT_HIGHLIGHT
tf_l = box_l.text_frame
tf_l.word_wrap = True
tf_l.margin_left = tf_l.margin_right = Inches(0.3)
tf_l.margin_top = Inches(0.3)

p_lt = tf_l.paragraphs[0]
p_lt.text = "Sự nóng chảy (Rắn ➔ Lỏng)"
p_lt.font.name = 'Arial'
p_lt.font.size = Pt(20)
p_lt.font.bold = True
p_lt.font.color.rgb = TEXT_TITLE
p_lt.space_after = Pt(14)

bullets_l = [
    "- Đun nóng chất rắn tinh khiết.",
    "- Phân tử thu nhiệt năng, dao động nhiệt tại chỗ mạnh lên không ngừng.",
    "- Biên độ dao động vượt qua giới hạn liên kết tinh thể cứng.",
    "- Mạng lưới cấu trúc tinh thể sụp đổ hoàn toàn.",
    "- Các phân tử bắt đầu trượt tự do lên nhau tạo thành trạng thái lỏng."
]
for item in bullets_l:
    p = tf_l.add_paragraph()
    p.text = item
    p.font.name = 'Arial'
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_BODY
    p.space_after = Pt(8)

# Cột Phải: Sự đông đặc
box_r = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.8))
box_r.fill.solid()
box_r.fill.fore_color.rgb = CARD_BG
box_r.line.color.rgb = TEXT_HIGHLIGHT
tf_r = box_r.text_frame
tf_r.word_wrap = True
tf_r.margin_left = tf_r.margin_right = Inches(0.3)
tf_r.margin_top = Inches(0.3)

p_rt = tf_r.paragraphs[0]
p_rt.text = "Sự đông đặc (Lỏng ➔ Rắn)"
p_rt.font.name = 'Arial'
p_rt.font.size = Pt(20)
p_rt.font.bold = True
p_rt.font.color.rgb = TEXT_TITLE
p_rt.space_after = Pt(14)

bullets_r = [
    "- Làm lạnh hoặc lấy nhiệt lượng ra khỏi chất lỏng.",
    "- Động năng chuyển động nhiệt phân tử giảm dần, hạt đi chậm lại.",
    "- Ở nhiệt độ đông đặc, lực tương tác phân tử bắt đầu kéo chúng sát lại.",
    "- Các phân tử bị bắt giữ và khóa chặt vào các nút mạng tinh thể mới.",
    "- Chất lỏng chuyển hóa thành chất rắn có cấu trúc cố định."
]
for item in bullets_r:
    p = tf_r.add_paragraph()
    p.text = item
    p.font.name = 'Arial'
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_BODY
    p.space_after = Pt(8)

# ==========================================
# SLIDE 13: GIẢI THÍCH SỰ CHUYỂN THỂ DƯỚI GÓC PHÂN TỬ - HÓA HƠI & NGƯNG TỤ
# ==========================================
slide13 = prs.slides.add_slide(slide_layout)
set_slide_background(slide13)
add_header(slide13, "CƠ CHẾ HÓA HƠI VÀ NGƯNG TỤ")

# Cột Trái: Sự hóa hơi
box_l = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.8))
box_l.fill.solid()
box_l.fill.fore_color.rgb = CARD_BG
box_l.line.color.rgb = TEXT_HIGHLIGHT
tf_l = box_l.text_frame
tf_l.word_wrap = True
tf_l.margin_left = tf_l.margin_right = Inches(0.3)
tf_l.margin_top = Inches(0.3)

p_lt = tf_l.paragraphs[0]
p_lt.text = "Sự hóa hơi (Lỏng ➔ Khí)"
p_lt.font.name = 'Arial'
p_lt.font.size = Pt(20)
p_lt.font.bold = True
p_lt.font.color.rgb = TEXT_TITLE
p_lt.space_after = Pt(14)

bullets_l = [
    "- Bay hơi: Các phân tử nước ở lớp bề mặt nhận động năng lớn, bứt phá khỏi lực hút phân tử xung quanh bay vào không khí (xảy ra ở mọi nhiệt độ).",
    "- Sôi: Tại nhiệt độ sôi, hóa hơi diễn ra dữ dội cả trên bề mặt lẫn sâu trong lòng chất lỏng dưới dạng các bọt khí nổi lên và vỡ ra."
]
for item in bullets_l:
    p = tf_l.add_paragraph()
    p.text = item
    p.font.name = 'Arial'
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_BODY
    p.space_after = Pt(8)
    p.line_spacing = 1.15

# Cột Phải: Sự ngưng tụ
box_r = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.8))
box_r.fill.solid()
box_r.fill.fore_color.rgb = CARD_BG
box_r.line.color.rgb = TEXT_HIGHLIGHT
tf_r = box_r.text_frame
tf_r.word_wrap = True
tf_r.margin_left = tf_r.margin_right = Inches(0.3)
tf_r.margin_top = Inches(0.3)

p_rt = tf_r.paragraphs[0]
p_rt.text = "Sự ngưng tụ (Khí ➔ Lỏng)"
p_rt.font.name = 'Arial'
p_rt.font.size = Pt(20)
p_rt.font.bold = True
p_rt.font.color.rgb = TEXT_TITLE
p_rt.space_after = Pt(14)

bullets_r = [
    "- Phân tử khí va chạm vào bề mặt vật thể có nhiệt độ lạnh.",
    "- Truyền bớt nhiệt năng ra môi trường ngoài, động năng giảm sút.",
    "- Chuyển động chậm lại, bị lực hút phân tử kéo xích lại gần nhau.",
    "- Các dải phân tử ngưng kết bám tụ thành các giọt chất lỏng lấp lánh."
]
for item in bullets_r:
    p = tf_r.add_paragraph()
    p.text = item
    p.font.name = 'Arial'
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_BODY
    p.space_after = Pt(8)

# ==========================================
# SLIDE 14: ỨNG DỤNG THỰC TẾ TRONG ĐỜI SỐNG VÀ CÔNG NGHIỆP
# ==========================================
slide14 = prs.slides.add_slide(slide_layout)
set_slide_background(slide14)
add_header(slide14, "ỨNG DỤNG TRONG THỰC TIỄN")

cols_app = [
    ("Luyện kim & Lò đúc", "Ứng dụng nóng chảy và đông đặc để biến các cục quặng kim loại thô thành thép lỏng đổ vào khuôn để đúc chi tiết máy, công cụ cơ khí hoàn chỉnh."),
    ("Sấy thăng hoa thực phẩm", "Thực phẩm được đông đá nhanh rồi đưa vào buồng hút chân không sâu. Nước đá thăng hoa trực tiếp thành khí giúp thực phẩm khô giòn, giữ 99% dinh dưỡng."),
    ("Tuần hoàn nước tự nhiên", "Mặt trời làm nước biển hóa hơi bay lên tạo mây, mây lạnh ngưng tụ thành hạt mưa rơi xuống, tuần hoàn duy trì điều hòa khí hậu mát mẻ cho toàn bộ địa cầu.")
]

for i, (title, desc) in enumerate(cols_app):
    x = left_margin + i * (col_width + gap)
    shape = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), col_width, Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = TEXT_HIGHLIGHT
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.space_after = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = 'Arial'
    p2.font.size = Pt(15)
    p2.font.color.rgb = TEXT_BODY
    p2.line_spacing = 1.25

# ==========================================
# SLIDE 15: CỦNG CỐ KIẾN THỨC & CÂU HỎI THỰC TẾ
# ==========================================
slide15 = prs.slides.add_slide(slide_layout)
set_slide_background(slide15)
add_header(slide15, "CỦNG CỐ & THẢO LUẬN")

# Tạo hộp câu hỏi
q_shape = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(1.6))
q_shape.fill.solid()
q_shape.fill.fore_color.rgb = CARD_BG
q_shape.line.color.rgb = TEXT_HIGHLIGHT
q_tf = q_shape.text_frame
q_tf.word_wrap = True
q_tf.margin_left = q_tf.margin_right = Inches(0.3)
q_tf.margin_top = Inches(0.2)

q_p = q_tf.paragraphs[0]
q_p.text = "CÂU HỎI TƯ DUY: Tại sao khi đun nóng một chất rắn tinh khiết đến nhiệt độ nóng chảy, nhiệt độ của nó không hề tăng lên trong suốt quá trình nóng chảy dù ta vẫn đun nóng liên tục?"
q_p.font.name = 'Arial'
q_p.font.size = Pt(18)
q_p.font.bold = True
q_p.font.color.rgb = TEXT_TITLE
q_p.line_spacing = 1.2

# Tạo hộp đáp án
a_shape = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.3), Inches(11.733), Inches(3.2))
a_shape.fill.solid()
a_shape.fill.fore_color.rgb = CARD_BG
a_shape.line.color.rgb = TEXT_TITLE
a_tf = a_shape.text_frame
a_tf.word_wrap = True
a_tf.margin_left = a_tf.margin_right = Inches(0.3)
a_tf.margin_top = Inches(0.3)

a_p = a_tf.paragraphs[0]
a_p.text = "GIẢI THÍCH BẢN CHẤT DƯỚI GÓC PHÂN TỬ:"
a_p.font.name = 'Arial'
a_p.font.size = Pt(18)
a_p.font.bold = True
a_p.font.color.rgb = TEXT_HIGHLIGHT
a_p.space_after = Pt(10)

answers = [
    "- Nhiệt độ đại diện cho động năng chuyển động nhiệt trung bình của các phân tử.",
    "- Khi chất rắn nóng chảy, toàn bộ nhiệt lượng thu vào được dùng để bẻ gãy, phá vỡ liên kết tinh thể cứng cáp (làm tăng thế năng tương tác phân tử).",
    "- Do động năng trung bình không đổi ➔ Nhiệt độ của chất được giữ cố định ở điểm nóng chảy cho đến khi chất rắn biến đổi hoàn toàn thành chất lỏng."
]

for item in answers:
    p = a_tf.add_paragraph()
    p.text = item
    p.font.name = 'Arial'
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_BODY
    p.space_after = Pt(8)

# ==========================================
# SLIDE 16: SLIDE KẾT THÚC
# ==========================================
slide16 = prs.slides.add_slide(slide_layout)
set_slide_background(slide16)

# Thêm hộp văn bản Tiêu đề
end_box = slide16.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
tf = end_box.text_frame
tf.word_wrap = True

# Lời cảm ơn
p1 = tf.paragraphs[0]
p1.text = "CẢM ƠN THẦY CÔ VÀ CÁC EM HỌC SINH!"
p1.alignment = PP_ALIGN.CENTER
p1.font.name = 'Arial'
p1.font.size = Pt(36)
p1.font.bold = True
p1.font.color.rgb = TEXT_TITLE
p1.space_after = Pt(20)

# Bài tập
p2 = tf.add_paragraph()
p2.text = "Giải đáp thắc mắc - Q&A  |  Bài tập về nhà: Câu 1, 2, 3 SGK trang 10"
p2.alignment = PP_ALIGN.CENTER
p2.font.name = 'Arial'
p2.font.size = Pt(18)
p2.font.color.rgb = TEXT_BODY
p2.space_after = Pt(30)

# Bài sau
p3 = tf.add_paragraph()
p3.text = "Chuẩn bị bài tiếp theo: Bài 2: Nội năng và Định luật I Nhiệt động lực học"
p3.alignment = PP_ALIGN.CENTER
p3.font.name = 'Arial'
p3.font.size = Pt(16)
p3.font.italic = True
p3.font.color.rgb = TEXT_HIGHLIGHT

# Lưu file slide PowerPoint (.pptx)
output_path = os.path.join(os.path.dirname(__file__), "Slides_PowerPoint.pptx")
prs.save(output_path)
print(f"Đã xuất file PowerPoint thành công tại: {output_path}")
